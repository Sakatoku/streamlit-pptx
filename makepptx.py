from pptx import  Presentation
import streamlit as st
import pandas as pd
import numpy as np
from PIL import Image

# このプロトタイプで使用する定数
OUTPUT_FILENAME = 'result.pptx'
OUTPUT_DIR = 'output'
TMP_DIR = 'tmp'
SLIDE_WIDTH = 16
SLIDE_HEIGHT = 9
EMU_PER_INCH = 914400

# デモ用のグラフの生成関数(1種類目)：積み上げバーチャート
@st.cache_data
def generate_example_figure1():
    import plotly.express as px # type: ignore
    # ランダムな値が入ったデータを生成
    years = np.arange(1970, 2024, 1)
    pref1 = ['Ehime', 7, 40]
    pref2 = ['Tokushima', 6, 30]
    pref3 = ['Kagawa', 5, 20]
    pref4 = ['Kochi', 5, 20]
    prefs = [pref1, pref2, pref3, pref4]
    df = pd.DataFrame()
    for pref in prefs:
        tmp_df = pd.DataFrame({'prefecture': pref[0], 'year': years, 'value': np.random.randint(pref[1], pref[2], len(years))})
        df = pd.concat([df, tmp_df])
    # 積み上げバーチャートを生成する
    fig = px.bar(df, x='year', y='value', color='prefecture',
                color_discrete_sequence=px.colors.qualitative.T10,
                barmode='stack', height=400)
    return ('積み上げバーチャート', fig)

# デモ用のグラフの生成関数(2種類目)：サンバーストチャート
@st.cache_data
def generate_example_figure2():
    import plotly.express as px # type: ignore
    # データを読み込む
    df = px.data.tips()
    # サンバーストチャートを生成する
    fig = px.sunburst(df, path=['sex', 'day', 'time'], values='total_bill', color='time',
                                      color_discrete_sequence=px.colors.qualitative.T10
        )
    return ('サンバーストチャート', fig)

# デモ用のグラフの生成関数(3種類目)：コロプレスマップ
@st.cache_data
def generate_example_figure3():
    import plotly.express as px # type: ignore
    from urllib.request import urlopen
    import json
    import pandas as pd
    # データを読み込む
    with urlopen('https://raw.githubusercontent.com/plotly/datasets/master/geojson-counties-fips.json') as response:
        counties = json.load(response)
        df = pd.read_csv("https://raw.githubusercontent.com/plotly/datasets/master/fips-unemp-16.csv", dtype={"fips": str})
        # コロプレスマップを生成する
        fig = px.choropleth_mapbox(df,
            geojson=counties,
            locations='fips',
            color='unemp',
            color_continuous_scale="Viridis",
            range_color=(0, 12),
            mapbox_style="carto-positron",
            zoom=3,
            center={"lat": 37.0902, "lon": -95.7129},
            opacity=0.5,
            labels={'unemp':'unemployment rate'}
        )
        fig.update_layout(margin={"r":0, "t":0, "l":0, "b":0})
    return  ('コロプレスマップ', fig)

# 画像を中央に貼り付けるためのサイズを計算する
def calculate_size(img, slide_size, rate):
    # 画像の縦横比を計算する
    img_size = Image.open(img).size
    img_aspect_ratio = img_size[0] / img_size[1]

    # スライドのサイズを取得する
    slide_width, slide_height = slide_size
    # スライドの縦横比を計算する
    slide_aspect_ratio = slide_width / slide_height

    # 画像の縦横比とスライドの縦横比を比較する
    if img_aspect_ratio > slide_aspect_ratio:
        # 画像のほうが横長の場合、スライドの幅に合わせて画像を縮小する
        plot_size = (slide_width * rate, slide_width / img_aspect_ratio * rate)
    else:
        # 画像のほうが縦長の場合、スライドの高さに合わせて画像を縮小する
        plot_size = (slide_height * img_aspect_ratio * rate, slide_height * rate)

    # スライドのサイズと調整後の画像のサイズから、画像を中央に配置するための座標を計算する
    x = (slide_width - plot_size[0]) / 2
    y = (slide_height - plot_size[1]) / 2
    return x, y, plot_size[0], plot_size[1]

# 画像を貼り付けたプレゼンテーションを作成する
def make_presentation(titles, imgs):
    # プレゼンテーションを新規作成する
    prs = Presentation('input/template.pptx')

    # レイアウトを選択。新規作成した場合は[6]が空スライド
    # 今回はtemplate.pptxを利用するため、[3]をハードコーディングする
    layout = prs.slide_layouts[3]

    # 画像ごとにスライドを作成する
    for i, img in enumerate(imgs):
        # スライドを追加する
        slide = prs.slides.add_slide(layout)
        # 追加したスライドからタイトルプレースホルダーを取得し、タイトルを設定する
        shapes = slide.shapes
        for shape in shapes:
            if shape.name.startswith('Title'):
                shape.text_frame.text = titles[i]
                break
        # 画像を中央に配置する
        size = calculate_size(img, [prs.slide_width, prs.slide_height], 0.7)
        _ = shapes.add_picture(img, size[0], size[1], size[2], size[3])
    
    return prs

# デモ用のグラフを生成する
title1, fig1 = generate_example_figure1()
title2, fig2 = generate_example_figure2()
title3, fig3 = generate_example_figure3()
# 生成したグラフを表示する
new_title1 = st.text_input('1st Chart', value=title1)
st.plotly_chart(fig1, use_container_width=True)
new_title2 = st.text_input('2nd Chart', value=title2)
st.plotly_chart(fig2, use_container_width=True)
new_title3 = st.text_input('3rd Chart', value=title3)
st.plotly_chart(fig3, use_container_width=True)

# 生成したグラフを画像として保存する
img_titlelist = [new_title1, new_title2, new_title3]
img_pathlist = []
for i, fig in enumerate([fig1, fig2, fig3]):
    img_path = f'{TMP_DIR}/fig{i}.png'
    # 画像化にはkaleidoというレンダリングエンジンの利用が推奨されている。ただし環境によってはフリーズすることが確認できたため、orcaを利用する
    # scaleを2にすることで、画像の解像度を上げることができる
    fig.write_image(img_path, format='png', scale=2, engine='orca')
    img_pathlist.append(img_path)

# 画像を中央に貼り付けたプレゼンテーションを作成する
prs = make_presentation(img_titlelist, img_pathlist)

# プレゼンテーションファイルを保存する
prs.save(f'{OUTPUT_DIR}/{OUTPUT_FILENAME}')

# 保存したプレゼンテーションファイルをダウンロードできるようにする
with open(f'{OUTPUT_DIR}/{OUTPUT_FILENAME}', 'rb') as f:
    bytes = f.read()
    st.download_button(
        label='Download PPTX',
        data=bytes,
        file_name=OUTPUT_FILENAME,
        mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
